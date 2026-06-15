"""
Generate Insurance Claims Processor slides and add them to the
Flogo-Industry-Vertical-UseCases.pptx presentation.
"""
import copy
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from lxml import etree

SRC = r"C:\Users\nshah\Downloads\Flogo-Industry-Vertical-UseCases.pptx"
OUT = r"C:\Users\nshah\Downloads\Flogo-Industry-Vertical-UseCases.pptx"

FONT = "Nata Sans"
DARK = RGBColor(0x24, 0x3C, 0x55)
PURPLE_DARK = RGBColor(0x36, 0x30, 0x73)
BLUE = RGBColor(0x33, 0x82, 0xD9)
ORANGE = RGBColor(0xFF, 0x61, 0x0A)
PURPLE = RGBColor(0x61, 0x38, 0xF0)
DARK_BLUE2 = RGBColor(0x25, 0x55, 0x81)
GRAY = RGBColor(0x66, 0x66, 0x66)
BODY = RGBColor(0x33, 0x33, 0x33)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY_BG = RGBColor(0xF5, 0xF5, 0xF5)
LIGHT_BLUE_BG = RGBColor(0xE3, 0xF1, 0xFF)
LIGHT_PURPLE_BG = RGBColor(0xEB, 0xE5, 0xFF)
GREEN = RGBColor(0x0B, 0x84, 0x57)
RED = RGBColor(0xDC, 0x26, 0x26)
AMBER = RGBColor(0xD9, 0x73, 0x06)
TEAL = RGBColor(0x0D, 0x96, 0x88)

# EMU sizes (1 inch = 914400 EMU, 1 pt = 12700 EMU)
LEFT_MARGIN = 640080
TOP_TITLE = 156948
TITLE_W = 7886700
TITLE_H = 519300
SUBTITLE_TOP = 713232
CONTENT_W = 7772400


def add_textbox(slide, left, top, width, height, text, font_size, bold=False,
                color=DARK, alignment=PP_ALIGN.LEFT, font_name=FONT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.size = font_size
    run.font.bold = bold
    run.font.name = font_name
    run.font.color.rgb = color
    return txBox


def add_multiline_textbox(slide, left, top, width, height, lines, font_size,
                          bold=False, color=DARK, alignment=PP_ALIGN.LEFT,
                          font_name=FONT, line_spacing=None):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line_text in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = alignment
        if line_spacing:
            p.space_after = line_spacing
        run = p.add_run()
        run.text = line_text
        run.font.size = font_size
        run.font.bold = bold
        run.font.name = font_name
        run.font.color.rgb = color
    return txBox


def add_rounded_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_rect(slide, left, top, width, height, fill_color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    return shape


def add_stat_card(slide, left, top, stat_value, stat_label, stat_detail,
                  accent_color, card_w=1874519, card_h=1051560):
    add_rounded_rect(slide, left, top, card_w, card_h, LIGHT_GRAY_BG)
    bar_h = 45720
    add_rect(slide, left, top, card_w, bar_h, accent_color)
    add_textbox(slide, left, top + 109728, card_w, Pt(36),
                stat_value, Pt(24), bold=True, color=accent_color,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + 594360, card_w, Pt(18),
                stat_label, Pt(10), bold=True, color=DARK,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, left, top + 804672, card_w, Pt(14.4),
                stat_detail, Pt(8), bold=False, color=GRAY,
                alignment=PP_ALIGN.CENTER)


def add_pain_point(slide, y_pos, label, description, dot_color=BLUE):
    add_rect(slide, LEFT_MARGIN, y_pos + 45720, 91440, 164592, dot_color)
    add_textbox(slide, LEFT_MARGIN + 164592, y_pos, 2926080, Pt(21.6),
                label, Pt(10), bold=True, color=DARK)
    add_textbox(slide, 3840480, y_pos, 4663440, Pt(21.6),
                description, Pt(9), bold=False, color=BODY)


def add_conversation_bubble(slide, left, top, width, height, role, text,
                            role_color, is_agent=False):
    bg_color = LIGHT_PURPLE_BG if is_agent else LIGHT_BLUE_BG
    add_rounded_rect(slide, left, top, width, height, bg_color)
    add_textbox(slide, left + 137160, top + 36576, 640200, Pt(12.6),
                role, Pt(8), bold=True, color=role_color)
    add_textbox(slide, left + 137160, top + 201168, width - 274320, height - 237744,
                text, Pt(8.5), bold=False, color=BODY)


def add_arch_tier(slide, top, label, label_color, boxes, box_w=None):
    tier_h = 502920
    add_rounded_rect(slide, LEFT_MARGIN, top, 7863840, tier_h, LIGHT_GRAY_BG)
    add_textbox(slide, LEFT_MARGIN + 137160, top, 1828800, tier_h,
                label, Pt(10), bold=True, color=label_color)
    if boxes:
        n = len(boxes)
        total_box_area = 7863840 - 2011680
        gap = 91440
        if box_w is None:
            box_w = (total_box_area - gap * (n - 1)) // n
        box_left = LEFT_MARGIN + 2011680
        box_top = top + (tier_h - 320040) // 2
        for i, box_text in enumerate(boxes):
            add_rounded_rect(slide, box_left, box_top, box_w, 320040,
                             RGBColor(0xE8, 0xE8, 0xE8))
            add_textbox(slide, box_left, box_top, box_w, 320040,
                        box_text, Pt(8), bold=True, color=DARK,
                        alignment=PP_ALIGN.CENTER)
            box_left += box_w + gap


def add_arrow(slide, top, width=183000):
    mid_x = LEFT_MARGIN + 7863840 // 2 - width // 2
    connector = slide.shapes.add_connector(
        1, mid_x, top, mid_x + width, top
    )
    connector.line.color.rgb = GRAY
    connector.line.width = Pt(1.5)


def add_step_flow(slide, x, y, number, title, detail, accent_color):
    circle_r = 182880
    add_rounded_rect(slide, x, y, circle_r, circle_r, accent_color)
    add_textbox(slide, x, y, circle_r, circle_r,
                number, Pt(12), bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + circle_r + 45720, 1100000, Pt(14),
                title, Pt(10), bold=True, color=DARK,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, x, y + circle_r + 228600, 1100000, Pt(26),
                detail, Pt(7.5), bold=False, color=GRAY,
                alignment=PP_ALIGN.CENTER)


def get_layout(prs, name):
    for layout in prs.slide_layouts:
        if layout.name == name:
            return layout
    return prs.slide_layouts[3]


def move_slide_to_end(prs, slide_index):
    """Move slide at slide_index to the end of the presentation."""
    nsmap = {'p': 'http://schemas.openxmlformats.org/presentationml/2006/main'}
    sldIdLst = prs.part._element.find('.//p:sldIdLst', nsmap)
    sldId_elements = list(sldIdLst)
    el = sldId_elements[slide_index]
    sldIdLst.remove(el)
    sldIdLst.append(el)


# ── Main ──────────────────────────────────────────────────────────────────

prs = Presentation(SRC)
end_slide_idx = len(prs.slides) - 1  # last slide (End Slide)

section_layout = get_layout(prs, "Section Breaker")
content_layout = get_layout(prs, "Overview / Summary")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 1: Section Breaker — Insurance Industry
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(section_layout)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:  # Title
        ph.text = "Insurance Industry"
        for run in ph.text_frame.paragraphs[0].runs:
            run.font.size = Pt(28)
            run.font.name = FONT
    elif ph.placeholder_format.idx == 1:  # Body
        tf = ph.text_frame
        p0 = tf.paragraphs[0]
        p0.text = "Insurance Claims AI Processor"
        for run in p0.runs:
            run.font.size = Pt(18)
            run.font.name = FONT
        p1 = tf.add_paragraph()
        run1 = p1.add_run()
        run1.text = "AI-Powered Claims Adjudication with Fraud Detection"
        run1.font.size = Pt(11)
        run1.font.name = FONT

# ════════════════════════════════════════════════════════════════════════
# SLIDE 2: The Problem — Claims Fraud
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(content_layout)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "Claims Processing — The $80B Problem"
        for run in ph.text_frame.paragraphs[0].runs:
            run.font.size = Pt(27)
            run.font.bold = False
            run.font.name = FONT

stat_top = 1020000
add_stat_card(slide, LEFT_MARGIN, stat_top,
              "$80B+", "Annual fraud losses", "across auto insurance in the US",
              BLUE)
add_stat_card(slide, LEFT_MARGIN + 1938527 * 1, stat_top,
              "30-45", "Days to resolve", "average claim processing time",
              ORANGE)
add_stat_card(slide, LEFT_MARGIN + 1938527 * 2, stat_top,
              "10-15%", "Claims are fraudulent", "industry-wide estimate",
              PURPLE)
add_stat_card(slide, LEFT_MARGIN + 1938527 * 3, stat_top,
              "60%", "Manual review rate", "claims requiring human intervention",
              DARK_BLUE2)

pain_y = 2240000
add_textbox(slide, LEFT_MARGIN, pain_y, CONTENT_W, Pt(20),
            "Pain points across claims operations", Pt(11),
            bold=True, color=DARK)

pain_y += 280000
add_pain_point(slide, pain_y,
               "Manual adjudication bottleneck",
               "Adjusters handle 150-200 open claims each — complex ones sit in queue for weeks")
add_pain_point(slide, pain_y + 274320,
               "Fraud detection is reactive",
               "SIU teams investigate after payout — by then money is gone and recovery is costly")
add_pain_point(slide, pain_y + 274320 * 2,
               "Siloed policy & fraud systems",
               "Adjusters toggle between policy admin, claims, and fraud databases to make a decision")
add_pain_point(slide, pain_y + 274320 * 3,
               "Inconsistent decisions",
               "Same claim type gets different outcomes depending on which adjuster reviews it")
add_pain_point(slide, pain_y + 274320 * 4,
               "No real-time risk scoring",
               "Fraud indicators are checked in batch overnight — not at time of submission")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 3: The Solution
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(content_layout)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "The Solution — Insurance Claims AI Processor"
        for run in ph.text_frame.paragraphs[0].runs:
            run.font.size = Pt(27)
            run.font.bold = False
            run.font.name = FONT

add_textbox(slide, LEFT_MARGIN, SUBTITLE_TOP, CONTENT_W, Pt(14.4),
            "A REST API that chains two LLM Client calls — policy verification via MCP + fraud detection via A2A",
            Pt(12), bold=False, color=GRAY)

# Solution cards
card_w = 3700000
card_h = 750000
card_gap = 400000
card_top = 1150000

# Card 1: Policy Verification
add_rounded_rect(slide, LEFT_MARGIN, card_top, card_w, card_h, LIGHT_BLUE_BG)
add_textbox(slide, LEFT_MARGIN + 137160, card_top + 68580, card_w - 274320, Pt(14),
            "Step 1: Policy Verification (MCP)", Pt(11), bold=True, color=BLUE)
add_multiline_textbox(slide, LEFT_MARGIN + 137160, card_top + 274320,
                      card_w - 274320, card_h - 342900,
                      ["LLM Client calls MCP Server tools to:",
                       "  •  Look up policy details (holder, vehicle, status)",
                       "  •  Check coverage for the specific claim type"],
                      Pt(8), color=BODY)

# Card 2: Fraud Assessment
add_rounded_rect(slide, LEFT_MARGIN + card_w + card_gap, card_top, card_w, card_h, LIGHT_PURPLE_BG)
add_textbox(slide, LEFT_MARGIN + card_w + card_gap + 137160, card_top + 68580,
            card_w - 274320, Pt(14),
            "Step 2: Fraud Assessment (A2A)", Pt(11), bold=True, color=PURPLE)
add_multiline_textbox(slide, LEFT_MARGIN + card_w + card_gap + 137160, card_top + 274320,
                      card_w - 274320, card_h - 342900,
                      ["LLM Client delegates to A2A agent to:",
                       "  •  Analyze claim patterns and suspicious indicators",
                       "  •  Calculate composite fraud risk score (0-100)"],
                      Pt(8), color=BODY)

# Arrow between cards
arrow_y = card_top + card_h // 2
arrow_left = LEFT_MARGIN + card_w
add_textbox(slide, arrow_left, arrow_y - Pt(8), card_gap, Pt(16),
            "→", Pt(18), bold=True, color=GRAY, alignment=PP_ALIGN.CENTER)

# Decision card
decision_top = card_top + card_h + 200000
decision_w = 7800000
add_rounded_rect(slide, LEFT_MARGIN, decision_top, decision_w, 500000,
                 RGBColor(0xF0, 0xFD, 0xF4))
add_textbox(slide, LEFT_MARGIN + 137160, decision_top + 68580, decision_w - 274320, Pt(14),
            "Final Decision: APPROVE  |  FLAG_FOR_REVIEW  |  DENY", Pt(11),
            bold=True, color=GREEN)
add_textbox(slide, LEFT_MARGIN + 137160, decision_top + 274320, decision_w - 274320, Pt(10),
            "LLM synthesizes policy coverage + fraud score into a single recommendation with reasoning",
            Pt(8.5), bold=False, color=BODY)

# Key differentiators
diff_top = decision_top + 680000
add_textbox(slide, LEFT_MARGIN, diff_top, CONTENT_W, Pt(16),
            "Key differentiator: LLM Client Activity", Pt(11), bold=True, color=DARK)

diff_items = [
    ("Dynamic LLM configuration", "Provider, model, API key set as activity inputs — switch at runtime"),
    ("Stateless one-shot inference", "No conversation memory needed — each claim is independent"),
    ("Sequential chaining", "Step 1 output feeds directly into Step 2 prompt — simple flow-level data passing"),
]
for i, (label, desc) in enumerate(diff_items):
    y = diff_top + 270000 + i * 250000
    add_rect(slide, LEFT_MARGIN, y + 45720, 91440, 140000, BLUE)
    add_textbox(slide, LEFT_MARGIN + 164592, y, 2400000, Pt(18),
                label, Pt(9), bold=True, color=DARK)
    add_textbox(slide, 3200000, y, 5300000, Pt(18),
                desc, Pt(8.5), bold=False, color=BODY)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 4: Architecture
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(content_layout)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "Reference Architecture — LLM Client + MCP + A2A"
        for run in ph.text_frame.paragraphs[0].runs:
            run.font.size = Pt(27)
            run.font.bold = False
            run.font.name = FONT

add_textbox(slide, LEFT_MARGIN, SUBTITLE_TOP, CONTENT_W, Pt(18),
            "Three independent Flogo microservices — lightweight Go binaries, deploy anywhere",
            Pt(9), bold=False, color=GRAY)

tier_gap = 140000
tier1_top = 980000
tier_h = 502920

# Tier 1: REST Client
add_arch_tier(slide, tier1_top, "REST Client", PURPLE_DARK,
              ["Postman", "curl", "Internal Portal", "Mobile App"])

add_arrow(slide, tier1_top + tier_h + 30000)

# Tier 2: Orchestrator
tier2_top = tier1_top + tier_h + tier_gap
add_rounded_rect(slide, LEFT_MARGIN, tier2_top, 7863840, 600000, LIGHT_BLUE_BG)
add_textbox(slide, LEFT_MARGIN + 137160, tier2_top, 1828800, 600000,
            "Orchestrator\n(port 9194)", Pt(10), bold=True, color=DARK)
# Sub-boxes inside orchestrator
orch_box_left = LEFT_MARGIN + 2011680
orch_box_top = tier2_top + 80000
orch_box_h = 200000
add_rounded_rect(slide, orch_box_left, orch_box_top, 2700000, orch_box_h,
                 RGBColor(0xD6, 0xEB, 0xFF))
add_textbox(slide, orch_box_left, orch_box_top, 2700000, orch_box_h,
            "LLM Client #1: LookupPolicy", Pt(8), bold=True, color=DARK,
            alignment=PP_ALIGN.CENTER)
add_textbox(slide, orch_box_left + 2700000, orch_box_top, 400000, orch_box_h,
            "→", Pt(14), bold=True, color=GRAY, alignment=PP_ALIGN.CENTER)
add_rounded_rect(slide, orch_box_left + 3100000, orch_box_top, 2700000, orch_box_h,
                 RGBColor(0xE5, 0xDB, 0xFF))
add_textbox(slide, orch_box_left + 3100000, orch_box_top, 2700000, orch_box_h,
            "LLM Client #2: AssessFraud", Pt(8), bold=True, color=DARK,
            alignment=PP_ALIGN.CENTER)
# LLM Provider label
add_textbox(slide, orch_box_left, orch_box_top + orch_box_h + 30000, 5800000, Pt(10),
            "LLM Provider: OpenAI / Anthropic / Gemini / Ollama  (dynamic config per step)",
            Pt(7), bold=False, color=GRAY, alignment=PP_ALIGN.CENTER)

# Arrows down to MCP and A2A
tier3_top = tier2_top + 600000 + tier_gap

# Tier 3a: MCP Server (left)
mcp_w = 3800000
add_rounded_rect(slide, LEFT_MARGIN, tier3_top, mcp_w, 650000, LIGHT_BLUE_BG)
add_textbox(slide, LEFT_MARGIN + 137160, tier3_top + 30000, mcp_w - 274320, Pt(12),
            "PolicyLookupMCPServer (port 9603)", Pt(9), bold=True, color=BLUE)
add_textbox(slide, LEFT_MARGIN + 137160, tier3_top + 220000, mcp_w - 274320, Pt(10),
            "MCP Server  •  Stateless  •  Read-only tools", Pt(7), bold=False, color=GRAY)
# Tool boxes
tool_top = tier3_top + 370000
tool_w = 1700000
tool_h = 230000
add_rounded_rect(slide, LEFT_MARGIN + 137160, tool_top, tool_w, tool_h,
                 RGBColor(0xD6, 0xEB, 0xFF))
add_textbox(slide, LEFT_MARGIN + 137160, tool_top, tool_w, tool_h,
            "lookup_policy", Pt(8), bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
add_rounded_rect(slide, LEFT_MARGIN + 137160 + tool_w + 100000, tool_top, tool_w, tool_h,
                 RGBColor(0xD6, 0xEB, 0xFF))
add_textbox(slide, LEFT_MARGIN + 137160 + tool_w + 100000, tool_top, tool_w, tool_h,
            "check_coverage", Pt(8), bold=True, color=DARK, alignment=PP_ALIGN.CENTER)

# Tier 3b: A2A Server (right)
a2a_left = LEFT_MARGIN + mcp_w + 263840
a2a_w = 7863840 - mcp_w - 263840
add_rounded_rect(slide, a2a_left, tier3_top, a2a_w, 650000, LIGHT_PURPLE_BG)
add_textbox(slide, a2a_left + 137160, tier3_top + 30000, a2a_w - 274320, Pt(12),
            "FraudDetectionA2A (port 9604)", Pt(9), bold=True, color=PURPLE)
add_textbox(slide, a2a_left + 137160, tier3_top + 220000, a2a_w - 274320, Pt(10),
            "A2A Server  •  Token Auth  •  PII Redaction", Pt(7), bold=False, color=GRAY)
# Tool boxes
add_rounded_rect(slide, a2a_left + 137160, tool_top, tool_w, tool_h,
                 RGBColor(0xE5, 0xDB, 0xFF))
add_textbox(slide, a2a_left + 137160, tool_top, tool_w, tool_h,
            "AnalyzePatterns", Pt(8), bold=True, color=DARK, alignment=PP_ALIGN.CENTER)
add_rounded_rect(slide, a2a_left + 137160 + tool_w + 100000, tool_top, tool_w, tool_h,
                 RGBColor(0xE5, 0xDB, 0xFF))
add_textbox(slide, a2a_left + 137160 + tool_w + 100000, tool_top, tool_w, tool_h,
            "CalculateRiskScore", Pt(8), bold=True, color=DARK, alignment=PP_ALIGN.CENTER)

# Bottom tier: Backend integration
backend_top = tier3_top + 650000 + tier_gap
add_arch_tier(slide, backend_top, "Backend Systems", PURPLE_DARK,
              ["Policy Admin DB", "Claims History", "Fraud Models (SAS/FICO)", "Underwriting API"])

# ════════════════════════════════════════════════════════════════════════
# SLIDE 5: How It Works — End-to-End Flow
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(content_layout)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "How It Works — End-to-End Flow"
        for run in ph.text_frame.paragraphs[0].runs:
            run.font.size = Pt(27)
            run.font.bold = False
            run.font.name = FONT

add_textbox(slide, LEFT_MARGIN, SUBTITLE_TOP, CONTENT_W, Pt(14.4),
            "One REST call triggers two sequential LLM Client steps — policy + fraud — returning a single decision",
            Pt(9), bold=False, color=GRAY)

steps = [
    ("1", "Claim submitted", "REST POST with\npolicy, type, amount,\nincident details", BLUE),
    ("2", "LLM Client #1", "System prompt:\n\"Verify policy and\ncheck coverage\"", RGBColor(0x29, 0x6F, 0xB2)),
    ("3", "lookup_policy", "MCP tool called\n→ holder, vehicle,\nstatus, history", DARK_BLUE2),
    ("4", "check_coverage", "MCP tool called\n→ COVERED /\nNOT_COVERED", PURPLE_DARK),
    ("5", "LLM Client #2", "System prompt:\n\"Assess fraud risk\nusing A2A agent\"", PURPLE),
    ("6", "AnalyzePatterns", "A2A tool called\n→ patterns,\nsuspicious flags", RGBColor(0x4D, 0x34, 0xB4)),
    ("7", "CalculateScore", "A2A tool called\n→ score 0-100,\nrecommendation", RGBColor(0x8B, 0x6D, 0xF4)),
    ("8", "Decision returned", "APPROVE /\nFLAG_FOR_REVIEW /\nDENY + reasoning", GREEN),
]

step_y = 1000000
step_x_start = 320000
step_spacing = 1080000
for i, (num, title, detail, color) in enumerate(steps):
    x = step_x_start + i * step_spacing
    add_step_flow(slide, x, step_y, num, title, detail, color)

# Flow arrow line
arrow_y_line = step_y + 91440
for i in range(7):
    x1 = step_x_start + i * step_spacing + 200000
    x2 = step_x_start + (i + 1) * step_spacing - 18000
    line = slide.shapes.add_connector(1, x1, arrow_y_line, x2, arrow_y_line)
    line.line.color.rgb = GRAY
    line.line.width = Pt(1)

# Annotations below the flow
annot_top = 3400000
annotations = [
    ("Steps 1-4: Policy Verification",
     "LLM Client #1 discovers and calls MCP tools automatically. The LLM decides which tools to call based on the system prompt and available tool schemas.",
     BLUE),
    ("Steps 5-7: Fraud Assessment",
     "LLM Client #2 delegates to the A2A fraud agent. The agent runs its own tools independently and returns structured fraud analysis.",
     PURPLE),
    ("Step 8: Decision Synthesis",
     "LLM Client #2 combines policy coverage results (from Step 1) with fraud scores (from Steps 5-7) into a final APPROVE / FLAG_FOR_REVIEW / DENY recommendation.",
     GREEN),
]

for i, (label, desc, color) in enumerate(annotations):
    y = annot_top + i * 350000
    add_rect(slide, LEFT_MARGIN, y + 30000, 91440, 164592, color)
    add_textbox(slide, LEFT_MARGIN + 164592, y, 2600000, Pt(21.6),
                label, Pt(9), bold=True, color=DARK)
    add_textbox(slide, 3100000, y, 5400000, Pt(21.6),
                desc, Pt(8), bold=False, color=BODY)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 6: Demo Scenarios — 5 Policies
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(content_layout)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "Demo Scenarios — 5 Built-In Policies"
        for run in ph.text_frame.paragraphs[0].runs:
            run.font.size = Pt(27)
            run.font.bold = False
            run.font.name = FONT

add_textbox(slide, LEFT_MARGIN, SUBTITLE_TOP, CONTENT_W, Pt(14.4),
            "Each policy produces a different outcome — from clean approval to high-fraud denial",
            Pt(9), bold=False, color=GRAY)

# Table-like rows
scenarios = [
    ("1", "POL-2026-001234", "James Morrison", "Honda CR-V 2025", "Collision $4,500",
     "15 LOW", "APPROVE", GREEN),
    ("2", "POL-2026-005678", "Sarah Mitchell", "BMW X5 2024", "Collision $18,500",
     "58 MED", "FLAG FOR REVIEW", AMBER),
    ("3", "POL-2026-009012", "Marcus Webb", "Mercedes GLE 2025", "Theft $62,000",
     "87 HIGH", "DENY", RED),
    ("4", "POL-2026-003456", "Elena Rodriguez", "Toyota RAV4 2023", "Medical $12,000",
     "22 LOW", "REVIEW (limit)", AMBER),
    ("5", "POL-2026-007890", "David Park", "Hyundai Tucson 2022", "Collision $8,500",
     "10", "DENY (expired)", RED),
]

# Column headers
header_top = 970000
col_positions = [LEFT_MARGIN, LEFT_MARGIN + 250000, LEFT_MARGIN + 1900000,
                 LEFT_MARGIN + 3400000, LEFT_MARGIN + 4900000,
                 LEFT_MARGIN + 6000000, LEFT_MARGIN + 6800000]
col_widths = [200000, 1600000, 1450000, 1450000, 1050000, 750000, 1200000]
headers = ["#", "Policy Number", "Holder", "Vehicle", "Claim", "Score", "Outcome"]

for j, (pos, w, hdr) in enumerate(zip(col_positions, col_widths, headers)):
    add_textbox(slide, pos, header_top, w, Pt(12),
                hdr, Pt(7.5), bold=True, color=DARK)

# Header underline
add_rect(slide, LEFT_MARGIN, header_top + 200000, 7863840, 15000, BLUE)

# Data rows
row_h = 320000
for i, (num, pol, holder, vehicle, claim, score, outcome, outcome_color) in enumerate(scenarios):
    y = header_top + 250000 + i * row_h
    # Alternating row background
    if i % 2 == 0:
        add_rounded_rect(slide, LEFT_MARGIN, y - 20000, 7863840, row_h - 20000,
                         RGBColor(0xFA, 0xFA, 0xFA))
    values = [num, pol, holder, vehicle, claim, score]
    for j, (pos, w, val) in enumerate(zip(col_positions, col_widths, values)):
        add_textbox(slide, pos, y, w, Pt(12),
                    val, Pt(7), bold=(j == 0), color=DARK if j > 0 else BLUE)
    # Outcome with color
    add_textbox(slide, col_positions[6], y, col_widths[6], Pt(12),
                outcome, Pt(7), bold=True, color=outcome_color)

# Scenario highlights
highlight_top = header_top + 250000 + 5 * row_h + 180000
highlights = [
    ("•  Scenario 1 (APPROVE): Clean 8-year customer, low-value claim, witness statement on file", GREEN),
    ("•  Scenario 2 (FLAG): Coverage limits increased 60 days before claim — suspicious timing", AMBER),
    ("•  Scenario 3 (DENY): New customer, 4 claims in 10 months (5x average), staged-loss pattern", RED),
    ("•  Scenario 4 (REVIEW): Fraud score is LOW but $12K medical exceeds $10K coverage limit", AMBER),
    ("•  Scenario 5 (DENY): Policy expired April 2026 — claim cannot be processed", RED),
]
for i, (text, color) in enumerate(highlights):
    add_textbox(slide, LEFT_MARGIN, highlight_top + i * 200000, 7863840, Pt(12),
                text, Pt(7.5), bold=False, color=color)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 7: Sample Request/Response — APPROVE vs DENY
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(content_layout)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "Sample Request & Response"
        for run in ph.text_frame.paragraphs[0].runs:
            run.font.size = Pt(27)
            run.font.bold = False
            run.font.name = FONT

add_textbox(slide, LEFT_MARGIN, SUBTITLE_TOP, CONTENT_W, Pt(14.4),
            "Side-by-side: a clean claim vs a high-fraud claim — same API, different outcomes",
            Pt(9), bold=False, color=GRAY)

# Left column: APPROVE
col_w = 3750000
col1_left = LEFT_MARGIN
col2_left = LEFT_MARGIN + col_w + 300000

# APPROVE header
approve_top = 950000
add_rounded_rect(slide, col1_left, approve_top, col_w, 280000,
                 RGBColor(0xDC, 0xFC, 0xE7))
add_textbox(slide, col1_left, approve_top, col_w, 280000,
            "Scenario 1: APPROVE", Pt(11), bold=True, color=GREEN,
            alignment=PP_ALIGN.CENTER)

# APPROVE request
req_top = approve_top + 350000
add_textbox(slide, col1_left, req_top, col_w, Pt(10),
            "Request:", Pt(7.5), bold=True, color=DARK)
add_rounded_rect(slide, col1_left, req_top + 150000, col_w, 550000,
                 RGBColor(0xF8, 0xF8, 0xF8))
add_multiline_textbox(slide, col1_left + 80000, req_top + 180000,
                      col_w - 160000, 500000,
                      ['policy_number: POL-2026-001234',
                       'claim_type: collision',
                       'claim_amount: $4,500',
                       'incident: "Vehicle struck in parking',
                       '  lot. Other driver left note."'],
                      Pt(6.5), color=BODY, font_name="Consolas")

# APPROVE response
resp_top = req_top + 750000
add_textbox(slide, col1_left, resp_top, col_w, Pt(10),
            "AI Response:", Pt(7.5), bold=True, color=DARK)
add_rounded_rect(slide, col1_left, resp_top + 150000, col_w, 1300000,
                 RGBColor(0xF0, 0xFD, 0xF4))
add_multiline_textbox(slide, col1_left + 80000, resp_top + 180000,
                      col_w - 160000, 1250000,
                      ['RECOMMENDATION: APPROVE',
                       '',
                       'Policy: ACTIVE (James Morrison)',
                       'Vehicle: 2025 Honda CR-V',
                       'Coverage: Collision COVERED, $50K limit',
                       'Deductible: $500',
                       '',
                       'Fraud Risk Score: 15/100 (LOW)',
                       '  Claim Frequency: 8 (1 in 5 years)',
                       '  Amount Anomaly: 12 (normal range)',
                       '',
                       'Mitigating: 8-year customer, safe driver',
                       'discount, third-party witness on file'],
                      Pt(6.5), color=GREEN, font_name="Consolas")

# DENY header
add_rounded_rect(slide, col2_left, approve_top, col_w, 280000,
                 RGBColor(0xFE, 0xE2, 0xE2))
add_textbox(slide, col2_left, approve_top, col_w, 280000,
            "Scenario 3: DENY", Pt(11), bold=True, color=RED,
            alignment=PP_ALIGN.CENTER)

# DENY request
add_textbox(slide, col2_left, req_top, col_w, Pt(10),
            "Request:", Pt(7.5), bold=True, color=DARK)
add_rounded_rect(slide, col2_left, req_top + 150000, col_w, 550000,
                 RGBColor(0xF8, 0xF8, 0xF8))
add_multiline_textbox(slide, col2_left + 80000, req_top + 180000,
                      col_w - 160000, 500000,
                      ['policy_number: POL-2026-009012',
                       'claim_type: theft',
                       'claim_amount: $62,000',
                       'incident: "Vehicle stolen from',
                       '  driveway overnight. No cameras."'],
                      Pt(6.5), color=BODY, font_name="Consolas")

# DENY response
add_textbox(slide, col2_left, resp_top, col_w, Pt(10),
            "AI Response:", Pt(7.5), bold=True, color=DARK)
add_rounded_rect(slide, col2_left, resp_top + 150000, col_w, 1300000,
                 RGBColor(0xFE, 0xF2, 0xF2))
add_multiline_textbox(slide, col2_left + 80000, resp_top + 180000,
                      col_w - 160000, 1250000,
                      ['RECOMMENDATION: DENY',
                       '',
                       'Policy: ACTIVE (Marcus Webb)',
                       'Vehicle: 2025 Mercedes GLE 450',
                       'Coverage: Theft COVERED, $50K limit',
                       'BUT claim $62K EXCEEDS limit',
                       '',
                       'Fraud Risk Score: 87/100 (HIGH)',
                       '  Claim Frequency: 95 (4 in 10 mo)',
                       '  Amount Anomaly: 92 (4.1x avg)',
                       '  Staged-loss pattern detected',
                       '',
                       'Risk: New customer, no evidence,',
                       'escalating claim amounts, no police rpt'],
                      Pt(6.5), color=RED, font_name="Consolas")

# ════════════════════════════════════════════════════════════════════════
# SLIDE 8: Extensibility
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(content_layout)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "Extensibility — Add New Capabilities in Minutes"
        for run in ph.text_frame.paragraphs[0].runs:
            run.font.size = Pt(27)
            run.font.bold = False
            run.font.name = FONT

add_textbox(slide, LEFT_MARGIN, SUBTITLE_TOP, CONTENT_W, Pt(14.4),
            "Extend the claims processor with new tools and integrations — no agent retraining required",
            Pt(9), bold=False, color=GRAY)

ext_items = [
    ("Claims History Lookup", "Pull full claims history from policy admin system for pattern analysis",
     BLUE),
    ("Document Verification", "AI-powered document consistency check — photos, police reports, receipts",
     PURPLE),
    ("SIU Referral", "Auto-escalate high-risk claims to Special Investigations Unit with full context",
     RED),
    ("Compliance Check", "Validate claim against state regulations and policy terms before decision",
     GREEN),
    ("Subrogation Analysis", "Identify recovery opportunities from third-party liable parties",
     DARK_BLUE2),
    ("Customer Communication", "Auto-generate status update emails and letters based on claim decision",
     ORANGE),
]

card_w_ext = 3700000
card_h_ext = 420000
card_gap_ext = 160000
cards_per_row = 2

for i, (title, desc, color) in enumerate(ext_items):
    row = i // cards_per_row
    col = i % cards_per_row
    x = LEFT_MARGIN + col * (card_w_ext + 400000)
    y = 980000 + row * (card_h_ext + card_gap_ext)
    add_rounded_rect(slide, x, y, card_w_ext, card_h_ext, LIGHT_GRAY_BG)
    add_rect(slide, x, y, 68580, card_h_ext, color)
    add_textbox(slide, x + 137160, y + 45720, card_w_ext - 205740, Pt(12),
                title, Pt(9), bold=True, color=DARK)
    add_textbox(slide, x + 137160, y + 228600, card_w_ext - 205740, card_h_ext - 274320,
                desc, Pt(7.5), bold=False, color=BODY)

# How to extend
ext_how_top = 980000 + 3 * (card_h_ext + card_gap_ext) + 100000
add_textbox(slide, LEFT_MARGIN, ext_how_top, CONTENT_W, Pt(16),
            "How easy is it to add a capability?", Pt(10), bold=True, color=DARK)
add_textbox(slide, LEFT_MARGIN, ext_how_top + 250000, CONTENT_W, Pt(14),
            "1. Build a Flogo flow that calls the new API   →   "
            "2. Register it as an MCP tool or A2A agent   →   "
            "3. Done. LLM Client discovers it automatically.",
            Pt(8.5), bold=False, color=BODY)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 9: Business Outcomes
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(content_layout)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "Business Outcomes"
        for run in ph.text_frame.paragraphs[0].runs:
            run.font.size = Pt(27)
            run.font.bold = False
            run.font.name = FONT

add_textbox(slide, LEFT_MARGIN, SUBTITLE_TOP, CONTENT_W, Pt(14.4),
            "Quantified value measurable from week 1 of production",
            Pt(9), bold=False, color=GRAY)

stat_top_bo = 1020000
add_stat_card(slide, LEFT_MARGIN, stat_top_bo,
              "< 30s", "Claim decision", "vs 30-45 days manual",
              BLUE)
add_stat_card(slide, LEFT_MARGIN + 1938527, stat_top_bo,
              "85%+", "Auto-adjudication", "clean claims processed without human touch",
              ORANGE)
add_stat_card(slide, LEFT_MARGIN + 1938527 * 2, stat_top_bo,
              "3x", "Fraud detection rate", "vs manual SIU review",
              PURPLE)
add_stat_card(slide, LEFT_MARGIN + 1938527 * 3, stat_top_bo,
              "$0", "New infrastructure", "lightweight Flogo microservices",
              DARK_BLUE2)

benefits_top = 2300000
add_textbox(slide, LEFT_MARGIN, benefits_top, CONTENT_W, Pt(20),
            "Additional benefits", Pt(11), bold=True, color=DARK)

benefits = [
    ("Consistent decisions",
     "Same claim type gets the same analysis regardless of which adjuster queue it lands in"),
    ("Real-time fraud scoring",
     "Risk assessment happens at submission time — not in overnight batch. Suspicious claims flagged immediately"),
    ("Regulatory compliance",
     "Full audit trail — every policy lookup, coverage check, and fraud score is logged with timestamps"),
    ("Scalable during catastrophes",
     "Handle surge volumes (storms, floods) without hiring temporary adjusters — AI scales horizontally"),
]

for i, (label, desc) in enumerate(benefits):
    y = benefits_top + 280000 + i * 274320
    add_rect(slide, LEFT_MARGIN, y + 45720, 91440, 164592, BLUE)
    add_textbox(slide, LEFT_MARGIN + 164592, y, 2926080, Pt(21.6),
                label, Pt(10), bold=True, color=DARK)
    add_textbox(slide, 3840480, y, 4663440, Pt(21.6),
                desc, Pt(9), bold=False, color=BODY)

# ════════════════════════════════════════════════════════════════════════
# SLIDE 10: Why TIBCO Flogo
# ════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(content_layout)
for ph in slide.placeholders:
    if ph.placeholder_format.idx == 0:
        ph.text = "Why TIBCO Flogo for Insurance Claims AI"
        for run in ph.text_frame.paragraphs[0].runs:
            run.font.size = Pt(27)
            run.font.bold = False
            run.font.name = FONT

why_items = [
    ("LLM Client Activity",
     "Dynamic LLM configuration — switch providers, models, and temperatures per step at runtime. No pre-configured connections needed.",
     BLUE),
    ("MCP + A2A in One Platform",
     "Policy lookup via MCP (stateless tools) + fraud detection via A2A (agent delegation) — both native Flogo capabilities in a single flow.",
     PURPLE),
    ("Sequential Chaining Pattern",
     "Step 1 output feeds directly into Step 2 prompt — simple flow-level data passing. No conversation memory or agent handoff infrastructure needed.",
     GREEN),
    ("Enterprise Security Built-In",
     "A2A Server with token authentication and PII redaction enabled by default. Sensitive policyholder data never leaves the secure boundary.",
     DARK_BLUE2),
    ("Lightweight & Deployable Anywhere",
     "Three independent Go microservices. Deploy on-prem behind the firewall, in the cloud, or at the edge. No JVM, no heavy middleware.",
     ORANGE),
    ("Observable & Auditable",
     "Every tool call logged with full context. OpenTelemetry support for end-to-end tracing. Complete audit trail for regulatory compliance.",
     RGBColor(0x4D, 0x34, 0xB4)),
]

why_card_w = 3700000
why_card_h = 500000
why_gap = 160000

for i, (title, desc, color) in enumerate(why_items):
    row = i // 2
    col = i % 2
    x = LEFT_MARGIN + col * (why_card_w + 400000)
    y = 880000 + row * (why_card_h + why_gap)
    add_rounded_rect(slide, x, y, why_card_w, why_card_h, LIGHT_GRAY_BG)
    add_rect(slide, x, y, 68580, why_card_h, color)
    add_textbox(slide, x + 137160, y + 45720, why_card_w - 205740, Pt(14),
                title, Pt(10), bold=True, color=color)
    add_textbox(slide, x + 137160, y + 250000, why_card_w - 205740, why_card_h - 296000,
                desc, Pt(7.5), bold=False, color=BODY)


# ── Reorder: move the original end slide to the very end ──────────────
move_slide_to_end(prs, end_slide_idx)

prs.save(OUT)
print(f"Saved to {OUT}")
print(f"Total slides: {len(prs.slides)}")
